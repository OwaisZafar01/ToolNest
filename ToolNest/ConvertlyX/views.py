
from django.shortcuts import render
from django.http import FileResponse, HttpResponse
from django.conf import settings
from django.core.files.storage import default_storage

import os
import pythoncom
import win32com.client

from docx2pdf import convert
from pdf2docx import Converter

from PIL import Image
from pptx import Presentation

import pandas as pd

from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from reportlab.lib.pagesizes import letter, landscape
from reportlab.platypus import (
    SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
)
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors

from xhtml2pdf import pisa

def index(request):
    return render(request, 'index.html')

def converter(request, conversion_type):
    return render(request, 'converter.html', {
        'conversion_type': conversion_type
    })

def upload(request, conversion_type):
    try:
        if request.method == 'POST':
            if conversion_type == 'word-to-pdf':
                return word_to_pdf_conversion(request)

            elif conversion_type == 'pdf-to-word':
                return pdf_to_word_conversion(request)

            elif conversion_type == 'image-to-pdf':
                return image_to_pdf_conversion(request)

            elif conversion_type == 'excel-to-pdf':
                return excel_to_pdf_conversion(request)

            elif conversion_type == 'txt-to-pdf':
                return txt_to_pdf_conversion(request)

            elif conversion_type == 'pptx-to-pdf':
                return pptx_to_pdf_conversion(request)

            elif conversion_type == 'html-to-pdf':
                return html_to_pdf_conversion(request)

            elif conversion_type == 'pdf-to-image':
                return pdf_to_image_conversion(request)
            
        return render(request, 'converter.html', {
            'conversion_type': conversion_type,
            'error': 'Invalid conversion method'
        })

    except Exception as e:
        return render(request, 'converter.html', {
            'conversion_type': conversion_type,
            'error': str(e)
        })

def word_to_pdf_conversion(request):
    if 'document' not in request.FILES:
        return render(request, 'converter.html', {
            'conversion_type': 'word-to-pdf',
            'error': 'No file uploaded'
        })

    uploaded_file = request.FILES['document']
    
    
    os.makedirs(os.path.join(settings.BASE_DIR, 'media', 'uploads'), exist_ok=True)
    
  
    file_path = os.path.join(settings.BASE_DIR, 'media', 'uploads', uploaded_file.name)
    with open(file_path, 'wb+') as destination:
        for chunk in uploaded_file.chunks():
            destination.write(chunk)
    
    pdf_path = file_path.replace('.docx', '.pdf')
    
    try:
        pythoncom.CoInitialize()
        word = win32com.client.Dispatch("Word.Application")
        doc = word.Documents.Open(file_path)
        doc.SaveAs(pdf_path, FileFormat=17)  
        doc.Close()
        word.Quit()
        pythoncom.CoUninitialize()
    except Exception as e:
        try:
            cv = Converter(file_path)
            cv.convert(pdf_path)
            cv.close()
        except Exception as fallback_error:
            return render(request, 'converter.html', {
                'conversion_type': 'word-to-pdf',
                'error': f'Conversion failed: {str(fallback_error)}'
            })
    
    if os.path.exists(pdf_path):
        response = FileResponse(
            open(pdf_path, 'rb'), 
            as_attachment=True, 
            filename=os.path.basename(pdf_path)
        )
        return response
    else:
        return render(request, 'converter.html', {
            'conversion_type': 'word-to-pdf',
            'error': 'PDF conversion failed'
        })

def pdf_to_word_conversion(request):

    if 'document' not in request.FILES:
        return render(request, 'converter.html', {
            'conversion_type': 'pdf-to-word',
            'error': 'No file uploaded'
        })

    uploaded_file = request.FILES['document']
    
    os.makedirs(os.path.join(settings.BASE_DIR, 'media', 'uploads'), exist_ok=True)
    
    file_path = os.path.join(settings.BASE_DIR, 'media', 'uploads', uploaded_file.name)
    with open(file_path, 'wb+') as destination:
        for chunk in uploaded_file.chunks():
            destination.write(chunk)
    
    word_path = file_path.replace('.pdf', '.docx')
    
    try:
        cv = Converter(file_path)
        cv.convert(word_path)
        cv.close()
    except Exception as e:
        return render(request, 'converter.html', {
            'conversion_type': 'pdf-to-word',
            'error': f'Conversion failed: {str(e)}'
        })
    
    if os.path.exists(word_path):
        response = FileResponse(
            open(word_path, 'rb'), 
            as_attachment=True, 
            filename=os.path.basename(word_path)
        )
        return response
    else:
        return render(request, 'converter.html', {
            'conversion_type': 'pdf-to-word',
            'error': 'PDF to Word conversion failed'
        })

def image_to_pdf_conversion(request):
    if 'images' not in request.FILES:
        return render(request, 'converter.html', {
            'conversion_type': 'image-to-pdf',
            'error': 'No images uploaded'
        })

    uploaded_images = request.FILES.getlist('images')
    
    os.makedirs(os.path.join(settings.BASE_DIR, 'media', 'uploads'), exist_ok=True)

    image_paths = []
    
    for uploaded_file in uploaded_images:
        file_path = os.path.join(settings.BASE_DIR, 'media', 'uploads', uploaded_file.name)
        with open(file_path, 'wb+') as destination:
            for chunk in uploaded_file.chunks():
                destination.write(chunk)
        image_paths.append(file_path)
    
    pdf_path = os.path.join(settings.BASE_DIR, 'media', 'uploads', 'converted_images.pdf')
    
    try:
        c = canvas.Canvas(pdf_path, pagesize=letter)
        width, height = letter
        
        for img_path in image_paths:
            img = Image.open(img_path)
            
            img_width, img_height = img.size
            aspect = img_height / img_width
            
            page_width = width - 2*inch
            scaled_width = page_width
            scaled_height = scaled_width * aspect
            
            if scaled_height > height - 2*inch:
                scaled_height = height - 2*inch
                scaled_width = scaled_height / aspect
            
            x_centered = (width - scaled_width) / 2
            y_centered = (height - scaled_height) / 2
            
            c.drawImage(img_path, x_centered, y_centered, width=scaled_width, height=scaled_height)
            c.showPage()
        
        c.save()
    except Exception as e:
        return render(request, 'converter.html', {
            'conversion_type': 'image-to-pdf',
            'error': f'Conversion failed: {str(e)}'
        })
    
    if os.path.exists(pdf_path):
        response = FileResponse(
            open(pdf_path, 'rb'), 
            as_attachment=True, 
            filename='converted_images.pdf'
        )
        return response
    else:
        return render(request, 'converter.html', {
            'conversion_type': 'image-to-pdf',
            'error': 'Image to PDF conversion failed'
        })


def excel_to_pdf_conversion(request):
    if 'document' not in request.FILES:
        return render(request, 'converter.html', {
            'conversion_type': 'excel-to-pdf',
            'error': 'No file uploaded'
        })

    uploaded_file = request.FILES['document']
    
    file_name = uploaded_file.name
    if not file_name.endswith(('.xlsx', '.xls')):
        return render(request, 'converter.html', {
            'conversion_type': 'excel-to-pdf',
            'error': 'Invalid file format. Please upload an Excel file (.xlsx or .xls)'
        })
    
    if uploaded_file.size > 10 * 1024 * 1024:
        return render(request, 'converter.html', {
            'conversion_type': 'excel-to-pdf',
            'error': 'File size exceeds the 10MB limit'
        })
    
    import uuid
    secure_filename = f"{uuid.uuid4().hex}_{os.path.basename(file_name)}"
    
    upload_dir = os.path.join(settings.BASE_DIR, 'media', 'uploads')
    os.makedirs(upload_dir, exist_ok=True)
    
    file_path = os.path.join(upload_dir, secure_filename)
    with open(file_path, 'wb+') as destination:
        for chunk in uploaded_file.chunks():
            destination.write(chunk)
    
    pdf_filename = os.path.splitext(secure_filename)[0] + '.pdf'
    pdf_path = os.path.join(upload_dir, pdf_filename)
    
    try:
        try:
            df = pd.read_excel(file_path, engine='openpyxl')
        except Exception as e1:
            try:
                df = pd.read_excel(file_path, engine='xlrd')
            except Exception as e2:
                os.remove(file_path)  # Clean up the uploaded file
                return render(request, 'converter.html', {
                    'conversion_type': 'excel-to-pdf',
                    'error': f'Failed to read Excel file: {str(e1)}. {str(e2)}'
                })
        

        if len(df) > 10000: 
            os.remove(file_path)
            return render(request, 'converter.html', {
                'conversion_type': 'excel-to-pdf',
                'error': 'File too large. Maximum 10,000 rows allowed.'
            })
        
        data = [df.columns.tolist()] + df.values.tolist()
        
        pdf = SimpleDocTemplate(
            pdf_path, 
            pagesize=landscape(letter),
            rightMargin=72, 
            leftMargin=72, 
            topMargin=72, 
            bottomMargin=18
        )
        
        table_style = TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.grey),
            ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
            ('ALIGN', (0,0), (-1,-1), 'CENTER'),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('FONTSIZE', (0,0), (-1,0), 12),
            ('BOTTOMPADDING', (0,0), (-1,0), 12),
            ('BACKGROUND', (0,1), (-1,-1), colors.beige),
            ('GRID', (0,0), (-1,-1), 1, colors.black)
        ])
        
        table = Table(data)
        table.setStyle(table_style)
        
        elements = [table]
        pdf.build(elements)

        os.remove(file_path)
    
    except Exception as e:
        if os.path.exists(file_path):
            os.remove(file_path)
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
            
        return render(request, 'converter.html', {
            'conversion_type': 'excel-to-pdf',
            'error': f'Conversion failed: {str(e)}'
        })
    
    if os.path.exists(pdf_path):
        response = FileResponse(
            open(pdf_path, 'rb'), 
            as_attachment=True, 
            filename=os.path.basename(uploaded_file.name).replace('.xlsx', '.pdf').replace('.xls', '.pdf')
        )
        
        def cleanup_pdf(response):
            if os.path.exists(pdf_path):
                os.remove(pdf_path)
            return response
        
        response.close = lambda: cleanup_pdf(response)
        return response
    else:
        return render(request, 'converter.html', {
            'conversion_type': 'excel-to-pdf',
            'error': 'Excel to PDF conversion failed'
        })

def txt_to_pdf_conversion(request):
    if 'document' not in request.FILES:
        return render(request, 'converter.html', {
            'conversion_type': 'txt-to-pdf',
            'error': 'No file uploaded'
        })

    uploaded_file = request.FILES['document']
    
   
    os.makedirs(os.path.join(settings.BASE_DIR, 'media', 'uploads'), exist_ok=True)
    
  
    file_path = os.path.join(settings.BASE_DIR, 'media', 'uploads', uploaded_file.name)
    with open(file_path, 'wb+') as destination:
        for chunk in uploaded_file.chunks():
            destination.write(chunk)
    
    pdf_path = file_path.replace('.txt', '.pdf')
    
    try:
        with open(file_path, 'r', encoding='utf-8') as text_file:
            text_content = text_file.read()
        
        pdf = SimpleDocTemplate(
            pdf_path, 
            pagesize=letter,
            rightMargin=72, 
            leftMargin=72, 
            topMargin=72, 
            bottomMargin=18
        )
        
        styles = getSampleStyleSheet()
        normal_style = styles['Normal']
        normal_style.fontSize = 11
        normal_style.leading = 14
    
        paragraphs = text_content.split('\n')
        
        elements = []
        for paragraph in paragraphs:
            if paragraph.strip():  
                para = Paragraph(paragraph, normal_style)
                elements.append(para)
                
                elements.append(Spacer(1, 6))
        
     
        pdf.build(elements)
    
    except Exception as e:
        return render(request, 'converter.html', {
            'conversion_type': 'txt-to-pdf',
            'error': f'Conversion failed: {str(e)}'
        })
    
    
    if os.path.exists(pdf_path):
        response = FileResponse(
            open(pdf_path, 'rb'), 
            as_attachment=True, 
            filename=os.path.basename(pdf_path)
        )
        return response
    else:
        return render(request, 'converter.html', {
            'conversion_type': 'txt-to-pdf',
            'error': 'Text to PDF conversion failed'
        })

def pptx_to_pdf_conversion(request):
    
    if 'document' not in request.FILES:
        return render(request, 'converter.html', {
            'conversion_type': 'pptx-to-pdf',
            'error': 'No file uploaded'
        })

    uploaded_file = request.FILES['document']
    
    
    os.makedirs(os.path.join(settings.BASE_DIR, 'media', 'uploads'), exist_ok=True)
    
    
    file_path = os.path.join(settings.BASE_DIR, 'media', 'uploads', uploaded_file.name)
    with open(file_path, 'wb+') as destination:
        for chunk in uploaded_file.chunks():
            destination.write(chunk)
    
    
    pdf_path = file_path.replace('.pptx', '.pdf').replace('.ppt', '.pdf')
    
    try:
        
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("Powerpoint.Application")
        
        presentation = powerpoint.Presentations.Open(file_path)
        
        presentation.SaveAs(pdf_path, 32) 
        
      
        presentation.Close()
        powerpoint.Quit()
        pythoncom.CoUninitialize()
    
    except Exception as win32_error:
        try:
            from pptx import Presentation
            from PIL import Image
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import inch
            
            prs = Presentation(file_path)
            
            c = canvas.Canvas(pdf_path)
            
            for slide in prs.slides:
              
                width, height = 8.5*inch, 11*inch
                c.setPageSize((width, height))
                
                for shape in slide.shapes:
                   
                    if shape.has_image:
                        image = shape.image
                        
                       
                        img_path = os.path.join(settings.BASE_DIR, 'media', 'uploads', f'temp_slide_{slide.slide_id}.png')
                        with open(img_path, 'wb') as f:
                            f.write(image.blob)
                        
                        c.drawImage(
                            img_path, 
                            width/4,  
                            height/4,  
                            width=width/2,
                            height=height/2
                        )
                
                c.showPage()
            
            c.save()
        
        except Exception as fallback_error:
            return render(request, 'converter.html', {
                'conversion_type': 'pptx-to-pdf',
                'error': f'Conversion failed: {str(fallback_error)}'
            })
    
    if os.path.exists(pdf_path):
        response = FileResponse(
            open(pdf_path, 'rb'), 
            as_attachment=True, 
            filename=os.path.basename(pdf_path)
        )
        return response
    else:
        return render(request, 'converter.html', {
            'conversion_type': 'pptx-to-pdf',
            'error': 'PowerPoint to PDF conversion failed'
        })


def html_to_pdf_conversion(request):
    if 'document' not in request.FILES:
        return render(request, 'converter.html', {
            'conversion_type': 'html-to-pdf',
            'error': 'No file uploaded'
        })

    uploaded_file = request.FILES['document']
    
    upload_dir = os.path.join(settings.BASE_DIR, 'media', 'uploads')
    os.makedirs(upload_dir, exist_ok=True)
    
    file_path = os.path.join(upload_dir, uploaded_file.name)
    with open(file_path, 'wb+') as destination:
        for chunk in uploaded_file.chunks():
            destination.write(chunk)
    
    pdf_path = file_path.replace('.html', '.pdf')
    
    try:
        with open(file_path, 'r', encoding='utf-8') as html_file:
            html_content = html_file.read()

        # Fix for xhtml2pdf compatibility
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <title>PDF Document</title>
        </head>
        <body>
        {html_content}
        </body>
        </html>
        """
        
        with open(pdf_path, 'wb') as pdf_file:
            pisa_status = pisa.CreatePDF(html_content, dest=pdf_file)
        
        if pisa_status.err:
            raise Exception("PDF conversion error")
    
    except Exception as conversion_error:
        return render(request, 'converter.html', {
            'conversion_type': 'html-to-pdf',
            'error': f'Conversion failed: {str(conversion_error)}'
        })
    
    if os.path.exists(pdf_path):
        response = FileResponse(
            open(pdf_path, 'rb'), 
            as_attachment=True, 
            filename=os.path.basename(pdf_path)
        )
        return response
    else:
        return render(request, 'converter.html', {
            'conversion_type': 'html-to-pdf',
            'error': 'HTML to PDF conversion failed'
        })


def pdf_to_image_conversion(request):
    if 'document' not in request.FILES:
        return render(request, 'converter.html', {
            'conversion_type': 'pdf-to-image',
            'error': 'No file uploaded'
        })

    uploaded_file = request.FILES['document']
    
    os.makedirs(os.path.join(settings.BASE_DIR, 'media', 'uploads'), exist_ok=True)
    os.makedirs(os.path.join(settings.BASE_DIR, 'media', 'output'), exist_ok=True)
    
    file_path = os.path.join(settings.BASE_DIR, 'media', 'uploads', uploaded_file.name)
    with open(file_path, 'wb+') as destination:
        for chunk in uploaded_file.chunks():
            destination.write(chunk)
    
    output_dir = os.path.join(settings.BASE_DIR, 'media', 'output')
    
    try:
       
        import fitz 
        
        pdf_document = fitz.open(file_path)
        image_files = []
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            pix = page.get_pixmap(alpha=False)
            output_filename = f"{os.path.splitext(uploaded_file.name)[0]}_page{page_num+1}.png"
            output_path = os.path.join(output_dir, output_filename)
            pix.save(output_path)
            image_files.append(output_path)
        
        pdf_document.close()
        response = FileResponse(open(image_files[0], 'rb'), as_attachment=True, filename=os.path.basename(image_files[0]))
        return response
    
    except ImportError:
        try:
            from pdf2image import convert_from_path
            
            images = convert_from_path(file_path)
            image_files = []
        
            for i, image in enumerate(images):
                output_filename = f"{os.path.splitext(uploaded_file.name)[0]}_page{i+1}.png"
                output_path = os.path.join(output_dir, output_filename)
                image.save(output_path, "PNG")
                image_files.append(output_path)
            
            response = FileResponse(open(image_files[0], 'rb'), as_attachment=True, filename=os.path.basename(image_files[0]))
            return response
        
        except Exception as fallback_error:
            return render(request, 'converter.html', {
                'conversion_type': 'pdf-to-image',
                'error': f'Conversion failed: {str(fallback_error)}'
            })
    
    except Exception as e:
        return render(request, 'converter.html', {
            'conversion_type': 'pdf-to-image',
            'error': f'Conversion failed: {str(e)}'
        })